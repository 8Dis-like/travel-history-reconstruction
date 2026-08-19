import sys, os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

pptx_path = r'g:\Abroad\UCLA\MSEng\Capstone\Securiport\docs\presentation\Presentation-03-master.pptx'
prs = pptx.Presentation(pptx_path)

WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(255, 209, 0)
LIGHT_GOLD = RGBColor(255, 229, 127)

slides_content = [
    {
        'title': 'Travel History Reconstruction from Travel Documents',
        'bullets': [
            'An End-to-End AI System Combining YOLOv26 Instance Segmentation & Visual Language Models',
            'Capstone Team: Hao Zhang, Zuyan Tao, Wilson Tee',
            'Faculty Advisor: Prof. Bruce Robins (UCLA Engineering)',
            'Industry Sponsor: Alvaro Ramirez (Securiport LLC)',
            'Presentation Date: August 25-26, 2026 | UCLA Samueli School of Engineering'
        ]
    },
    {
        'title': 'Meet the Team',
        'bullets': [
            'Hao Zhang - Detection & Integration Lead: Preprocessing, Multi-Rotation Strategy, YOLO training, System integration, Report authoring.',
            'Zuyan Tao - VLM/OCR & Backend Lead: VLM integration, FastAPI backend, OCR pipeline, Timeline reconstruction algorithm.',
            'Wilson Tee - Dataset & Frontend Lead: Synthetic dataset pipeline, YOLOv26 segmentation model, React + Ant Design web UI.',
            'Advisors: Prof. Bruce Robins (UCLA Faculty Advisor) & Alvaro Ramirez (Securiport LLC Industry Sponsor).'
        ]
    },
    {
        'title': 'Presentation Agenda',
        'bullets': [
            '01. Problem Statement & Motivation',
            '02. Project Objectives & Scope',
            '03. Background & Related Work',
            '04. System Architecture (5-Stage Pipeline)',
            '05. Dataset Construction (U-Net & Overlap Bias)',
            '06. Detection Model Results (YOLOv8 vs YOLOv26)',
            '07. VLM Extraction Benchmarks & Rotation Robustness',
            '08. System-Level Evaluation (5 Real Passports)',
            '09. Live Web Application Demo (React + FastAPI)',
            '10. Conclusions & Future Work'
        ]
    },
    {
        'title': 'The Problem Statement',
        'bullets': [
            'Processing Bottleneck: Passport stamps are placed non-chronologically across multiple pages. Manual cross-referencing is extremely slow.',
            'Diversity & Degradation: 190+ countries, faded ink, blurred seals, overlapping stamps, and arbitrary rotation angles degrade classic OCR.',
            'Industry Need (Securiport): Automated document intelligence solution to extract entry/exit events and present a reviewable timeline.',
            'Research Gap: Prior work (Zaaboub et al. 2020) relies on rigid classifiers that fail on real-world degraded & overlapping stamps.'
        ]
    },
    {
        'title': 'Objectives & Scope',
        'bullets': [
            '01. Synthetic Dataset: Generated 5,000 synthetic training images using 871 stamps recovered via U-Net segmentation.',
            '02. Instance Segmentation Model: Trained YOLOv26s-seg with 4-rotation strategy for orientation-robust stamp localization.',
            '03. VLM Extraction Layer: Config-driven VLM layer extracting structured JSON (date, country, entry/exit direction).',
            '04. Human-in-the-Loop Web App: React + FastAPI web dashboard with interactive timeline table, SVG overlays, and edit loop.'
        ]
    },
    {
        'title': 'Background & Related Work',
        'bullets': [
            'Object Detection & Document AI: YOLO family (Redmon 2016 -> YOLOv26 2025) for real-time localization. LayoutLM & Donut (Xu 2020, Kim 2022) for VDU.',
            'Limitations of Traditional Pipelines: Multistage OCR (Tesseract, Paddle) degrades on rotated, ink-faded, and overlapping stamp images.',
            'Visual Language Models: Claude 3.5 Sonnet, Gemini 1.5 Flash, Qwen2-VL replace brittle OCR engines with joint vision-text reasoning.',
            'Our Novel Approach: Combine YOLOv26 instance segmentation (stamp localization) with config-driven VLM field extraction.'
        ]
    },
    {
        'title': 'System Architecture',
        'bullets': [
            'Stage 1 - Upload & Rasterization: PyMuPDF splits PDF pages into 300DPI rasters and 96DPI thumbnails.',
            'Stage 2 - Stamp Detection: YOLOv26s-seg performs 4-rotation inference (0, 90, 180, 270 deg) to extract polygon masks.',
            'Stage 3 - VLM Extraction: Crop regions fed to BaseExtractor (Claude / Gemini API) returning structured JSON fields.',
            'Stage 4 - Timeline Reconstruction: Python pairing algorithm matches entry/exit events and flags conflicts.',
            'Stage 5 - Web Dashboard: React + Ant Design interface allows analysts to verify, edit, and rebuild timelines in real time.'
        ]
    },
    {
        'title': 'Synthetic Dataset Construction',
        'bullets': [
            'Raw Stamp Scraping: 1,517 raw stamp images collected from Wikimedia Commons.',
            'U-Net Segmentation: ResNet-34 U-Net recovered 490 stamps beyond Otsu thresholding - doubling stamp library to 871 stamps.',
            'Synthetic Compositing: 5,000 synthetic pages (4,000 train / 1,000 val) generated with realistic overlap bias.',
            'Real-World Test Set: 194 real passport pages containing 1,002 hand-labeled stamps for rigorous benchmark evaluation.'
        ]
    },
    {
        'title': 'Stamp Detection Results',
        'bullets': [
            'Synthetic Validation: YOLOv8s achieved mAP50 of 0.9950, mAP50-95 of 0.9784, Precision 0.9983, Recall 0.9968.',
            'Real-World Baseline (YOLOv8s Box): Real-world recall was 83.3% with an F1 score of 0.859.',
            'Real-World Upgrade (YOLOv26s-seg): Instance segmentation raised real-world recall to 94.8% (+11.5%) and F1 to 0.927.',
            'Key Finding: Polygon segmentation masks yield tighter crops, eliminating cross-stamp bounding box contamination.'
        ]
    },
    {
        'title': 'VLM Extraction Benchmarks',
        'bullets': [
            'Field Extraction Accuracy (30 Curated Stamps): Claude 3.5 Sonnet (100%), Gemini 1.5 Flash (100%), Qwen2-VL-2B (~40%).',
            'Rotation Robustness (Claude Sonnet, 7 Angles): 0-75 deg (High accuracy), 90 deg (Moderate date, High country/direction), 180-220 deg (Degraded).',
            'Local VLM Bottleneck: Qwen2-VL-2B (4.5GB VRAM) achieved only 18.2% date accuracy and 45.5% country accuracy on degraded crops.',
            'Recommendation: Use Claude API as primary backend; config-driven factory enables seamless local model fallback.'
        ]
    },
    {
        'title': 'System-Level Pipeline Evaluation',
        'bullets': [
            '5 Full Real Passport Benchmark: Evaluated end-to-end timeline reconstruction accuracy across complete travel histories.',
            'Cascading Error Principle: Stage 1 Detection Recall (94.8%) -> Stage 2 Stamp Extraction (62.3%) -> Stage 3 Timeline Reconstruction (35.8%).',
            'Core Capstone Finding: The ~27-point gap between stamp extraction (62.3%) and timeline accuracy (35.8%) proves reconstruction logic is primary bottleneck.',
            'Stay Correctness Criterion: A stay is correct if and only if all stamps detected, fields extracted, entry/exit paired, and chronological order kept.'
        ]
    },
    {
        'title': 'Live Web Application Demo',
        'bullets': [
            'Multi-Page Upload: Server-side PDF rendering with thumbnail data URLs for low browser memory pressure.',
            'Interactive Overlays: SVG polygon masks overlaid directly on passport page scans with hover highlight.',
            'Status-Flagged Timeline: Color-coded stays (Confirmed, Inferred, Flagged Missing Exit / Conflicting Country).',
            'Human-in-the-Loop Review: One-click field editing and instant timeline rebuild loop for border analysts.'
        ]
    },
    {
        'title': 'Conclusions & Major Contributions',
        'bullets': [
            '1. U-Net & YOLOv26 Upgrade: Doubled stamp dataset to 871 stamps. Achieved 94.8% Recall and 0.927 F1 on 1,002 real test stamps.',
            '2. VLM-Over-OCR Validation: Empirical evidence that multimodal VLMs significantly outperform 4-stage OCR on degraded stamps.',
            '3. Triage-First Dashboard: Delivered React/FastAPI review system pairing every extracted stay with its source page scan.',
            'Deployment Recommendation: Deploy as a human-in-the-loop triage tool for border control officers.'
        ]
    },
    {
        'title': 'Future Work & Acknowledgements',
        'bullets': [
            'Advanced Conflict Resolution: Geographic plausibility checks and date inferencing to close the 27-point timeline gap.',
            'Contamination Isolation: Dedicated experiment isolating cross-stamp crop contamination failure modes.',
            'Multilingual Expansion: Support for non-Latin stamp scripts (Arabic, Cyrillic, CJK).',
            'Acknowledgements: Special thanks to Prof. Bruce Robins (UCLA Advisor) and Alvaro Ramirez (Securiport LLC Sponsor).'
        ]
    }
]

for idx, data in enumerate(slides_content):
    if idx < len(prs.slides):
        slide = prs.slides[idx]
        
        # Set Title Placeholder
        title_shape = None
        try:
            title_shape = slide.shapes.title
        except:
            pass
            
        if not title_shape:
            for s in slide.shapes:
                if s.has_text_frame and ('title' in s.text.lower() or 'header' in s.text.lower() or 'travel' in s.text.lower()):
                    title_shape = s
                    break
        
        if title_shape and title_shape.has_text_frame:
            tf = title_shape.text_frame
            tf.text = data['title']
            for p in tf.paragraphs:
                p.font.color.rgb = WHITE
                p.font.bold = True
                p.font.size = Pt(22 if idx > 0 else 32)

        # Set Content Body Placeholder
        body_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape != title_shape:
                txt = shape.text_frame.text.strip().lower()
                if txt == '' or 'click' in txt or 'department' in txt or 'label' in txt or '00%' in txt or 'synthetic' in txt or 'stage' in txt or 'object' in txt:
                    body_shape = shape
                    break
        
        if not body_shape:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8))
            tf = txBox.text_frame
        else:
            tf = body_shape.text_frame
        
        tf.word_wrap = True
        tf.text = '' # Clear
        
        for b_idx, bullet in enumerate(data['bullets']):
            p = tf.add_paragraph() if b_idx > 0 or len(tf.paragraphs[0].text) > 0 else tf.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(13 if idx > 0 else 15)
            p.font.color.rgb = WHITE
            p.space_after = Pt(10)

prs.save(r'g:\Abroad\UCLA\MSEng\Capstone\Securiport\docs\presentation\Presentation-03-master.pptx')
print('SUCCESSFULLY updated Presentation-03-master.pptx with clean headers and body bullet frames!')
